import base64
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any

from pydantic import Field, field_validator

from ..models import StrictModel
from ..source_types import (
    MediaAttachment,
    SourceCapabilityError,
    SourceCommand,
    SourceContext,
    SourceDataError,
)

RENDER_TIMEOUT_SECONDS = 60
MAX_RENDERED_SLIDES = 60
MAX_RENDERED_BYTES = 24 * 1024 * 1024


class PptxExtractTextInput(StrictModel):
    """Input for pptx.extract_text.

    Attributes:
        path: Workspace-relative PPTX path.
    """

    path: str

    @field_validator("path")
    @classmethod
    def require_pptx_extension(cls, value: str) -> str:
        """Requires the path to target a PPTX file.

        Args:
            value: Workspace-relative path from verifier.json.

        Returns:
            The unchanged path when it has a PPTX extension.

        Raises:
            ValueError: If the path does not end in ``.pptx``.
        """
        if Path(value).suffix.lower() != ".pptx":
            raise ValueError("pptx.extract_text requires a .pptx file")
        return value


class PptxExtractTextOutput(StrictModel):
    """Output from pptx.extract_text.

    Attributes:
        text: Extracted slide text.
    """

    text: str


class ExtractText(SourceCommand[PptxExtractTextInput, PptxExtractTextOutput]):
    """Extracts text from PPTX slides."""

    name = "extract_text"
    input_model = PptxExtractTextInput
    output_model = PptxExtractTextOutput

    def run(self, source_input: PptxExtractTextInput, context: SourceContext) -> PptxExtractTextOutput:
        """Runs PPTX text extraction.

        Args:
            source_input: Validated PPTX extraction input.
            context: Source runtime context.

        Returns:
            Extracted PPTX text.
        """
        resolved = context.resolve_path(source_input.path)
        return PptxExtractTextOutput(text=extract_pptx_text(resolved, context.max_content_chars))


class PptxRenderSlidesInput(StrictModel):
    """Input for ``pptx.render_slides``."""

    path: str

    @field_validator("path")
    @classmethod
    def require_pptx_extension(cls, value: str) -> str:
        if Path(value).suffix.lower() != ".pptx":
            raise ValueError("pptx.render_slides requires a .pptx file")
        return value


class RenderedSlide(StrictModel):
    """Visual evidence for one slide, with pixels excluded from JSON output."""

    number: int
    width: int
    height: int
    image_base64: str = Field(exclude=True, repr=False)


class PptxRenderSlidesOutput(StrictModel):
    """Rendered slide evidence plus ordinary extracted text."""

    text: str
    slide_count: int
    slides: list[RenderedSlide]
    truncated: bool
    renderer: str

    def media_attachments(self) -> list[MediaAttachment]:
        """Return image payloads without exposing them to reward serialization."""
        return [
            MediaAttachment(
                label=f"Slide {slide.number}",
                mime_type="image/png",
                data_base64=slide.image_base64,
            )
            for slide in self.slides
        ]


class RenderSlides(SourceCommand[PptxRenderSlidesInput, PptxRenderSlidesOutput]):
    """Render slide images for a prompt-grounded visual rubric."""

    name = "render_slides"
    input_model = PptxRenderSlidesInput
    output_model = PptxRenderSlidesOutput

    def run(
        self, source_input: PptxRenderSlidesInput, context: SourceContext
    ) -> PptxRenderSlidesOutput:
        resolved = context.resolve_path(source_input.path)
        return render_pptx_slides(resolved, context.max_content_chars)


class PptxInspectDeckInput(StrictModel):
    """Input for ``pptx.inspect_deck``.

    Attributes:
        path: Workspace-relative PPTX path.
    """

    path: str

    @field_validator("path")
    @classmethod
    def require_pptx_extension(cls, value: str) -> str:
        """Requires the path to target a PPTX file.

        Args:
            value: Workspace-relative path from verifier.json.

        Returns:
            The unchanged path when it has a PPTX extension.

        Raises:
            ValueError: If the path does not end in ``.pptx``.
        """
        if Path(value).suffix.lower() != ".pptx":
            raise ValueError("pptx.inspect_deck requires a .pptx file")
        return value


class PptxInspectDeckOutput(StrictModel):
    """Deck structure expressed as counts, names, and speaker notes.

    python-pptx only: unlike ``pptx.render_slides`` this needs no LibreOffice,
    produces no images, and carries no image-capability requirement, so it can
    answer "does the deck have four slides" deterministically and for free.
    It sees structure, never visual quality; a presentation-design obligation
    still belongs to ``pptx.render_slides`` and a rubric.

    Attributes:
        slide_count: Slides stored in the deck, hidden ones included.
        visible_slide_count: Slides an audience is actually shown. PowerPoint
            hides a slide with ``show="0"``, which a stored count cannot see, so
            a prompt capping what the audience sees means this number.
        hidden_slide_count: Slides marked hidden, such as a backup appendix.
        slide_titles: One entry per slide, in deck order, empty for a slide
            whose layout has no title placeholder or whose title is blank.
            Aligning with ``slide_count`` keeps a position meaningful and makes
            ``not_contains ""`` mean "every slide is titled".
        layout_names: One slide-layout name per slide, in deck order.
        image_count: Picture shapes across the deck.
        chart_count: Chart shapes across the deck.
        table_count: Table shapes across the deck.
        slides_with_notes: Slides carrying non-empty speaker notes. Exact even
            when ``notes_text`` is truncated.
        notes_text: Speaker notes, each prefixed with its slide number and
            capped to the configured source content limit.
    """

    slide_count: int
    visible_slide_count: int
    hidden_slide_count: int
    slide_titles: list[str]
    layout_names: list[str]
    image_count: int
    chart_count: int
    table_count: int
    slides_with_notes: int
    notes_text: str


class InspectDeck(SourceCommand[PptxInspectDeckInput, PptxInspectDeckOutput]):
    """Report slide count, titles, layouts, embedded objects, and notes."""

    name = "inspect_deck"
    input_model = PptxInspectDeckInput
    output_model = PptxInspectDeckOutput

    def run(
        self, source_input: PptxInspectDeckInput, context: SourceContext
    ) -> PptxInspectDeckOutput:
        """Runs deck structure inspection.

        Args:
            source_input: Validated deck inspection input.
            context: Source runtime context.

        Returns:
            Deck structure as counts, names, and bounded notes text.
        """
        resolved = context.resolve_path(source_input.path)
        return inspect_pptx_deck(resolved, context.max_content_chars)


def inspect_pptx_deck(path: Path, limit: int) -> PptxInspectDeckOutput:
    """Reads deck structure without rendering anything.

    Args:
        path: Resolved PPTX file path.
        limit: Maximum number of speaker-notes characters to return.

    Returns:
        Deck structure as counts, names, and notes text capped to ``limit``.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    picture_types = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}
    presentation = Presentation(str(path))
    titles: list[str] = []
    layout_names: list[str] = []
    notes_parts: list[str] = []
    notes_total = 0
    image_count = 0
    chart_count = 0
    table_count = 0
    slides_with_notes = 0
    hidden_count = 0

    for index, slide in enumerate(presentation.slides, start=1):
        if slide_is_hidden(slide):
            hidden_count += 1
        title_shape = slide.shapes.title
        titles.append("" if title_shape is None else (title_shape.text or "").strip())
        layout_names.append(str(getattr(slide.slide_layout, "name", "") or ""))
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                chart_count += 1
            if getattr(shape, "has_table", False):
                table_count += 1
            if getattr(shape, "shape_type", None) in picture_types:
                image_count += 1
        # `slide.notes_slide` creates the notes part on access, so the guard has
        # to be `has_notes_slide` rather than a try/except around the read.
        if not slide.has_notes_slide:
            continue
        notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if not notes:
            continue
        # Counted before the content bound: a truncated evidence string must
        # not silently lower the number an author asserts on.
        slides_with_notes += 1
        if notes_total < limit:
            notes_total = append_limited(
                notes_parts, f"[Slide {index}] {notes}", notes_total, limit
            )

    return PptxInspectDeckOutput(
        slide_count=len(titles),
        visible_slide_count=len(titles) - hidden_count,
        hidden_slide_count=hidden_count,
        slide_titles=titles,
        layout_names=layout_names,
        image_count=image_count,
        chart_count=chart_count,
        table_count=table_count,
        slides_with_notes=slides_with_notes,
        notes_text="\n".join(notes_parts)[:limit],
    )


def extract_pptx_text(path: Path, limit: int) -> str:
    """Extracts visible slide and table text from a PPTX file.

    Args:
        path: Resolved PPTX file path.
        limit: Maximum number of characters to return.

    Returns:
        Extracted text capped to ``limit`` characters.
    """
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    total = 0
    for index, slide in enumerate(presentation.slides, start=1):
        slide_parts = list(iter_slide_text(slide))
        if not slide_parts:
            continue
        total = append_limited(parts, f"[Slide {index}]", total, limit)
        if total >= limit:
            break
        for text in slide_parts:
            total = append_limited(parts, text, total, limit)
            if total >= limit:
                break
        if total >= limit:
            break
    return "\n".join(parts)[:limit]


def render_pptx_slides(path: Path, text_limit: int) -> PptxRenderSlidesOutput:
    """Render a PPTX through isolated LibreOffice and Poppler subprocesses.

    Rendering happens only in a temporary directory. The submitted deck is
    never mutated, macros are never run, and the LibreOffice user profile is
    isolated so parallel benchmark runs cannot share state.
    """
    # Parse before invoking LibreOffice so a malformed ZIP becomes the same
    # normal agent-output failure as pptx.extract_text.
    text = extract_pptx_text(path, text_limit)
    office = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not office or not pdftoppm:
        raise SourceCapabilityError(
            "pptx.render_slides requires LibreOffice Impress and pdftoppm in the benchmark image"
        )

    with tempfile.TemporaryDirectory(prefix="filecheck-pptx-") as raw_temp:
        tempdir = Path(raw_temp)
        profile = tempdir / "profile"
        pdf_dir = tempdir / "pdf"
        png_dir = tempdir / "png"
        pdf_dir.mkdir()
        png_dir.mkdir()
        try:
            subprocess.run(
                [
                    office,
                    "--headless",
                    "--nologo",
                    "--nodefault",
                    "--nofirststartwizard",
                    f"-env:UserInstallation={profile.as_uri()}",
                    "--convert-to",
                    "pdf:impress_pdf_Export",
                    "--outdir",
                    str(pdf_dir),
                    str(path),
                ],
                check=True,
                capture_output=True,
                text=True,
                timeout=RENDER_TIMEOUT_SECONDS,
            )
            pdf_files = sorted(pdf_dir.glob("*.pdf"))
            if len(pdf_files) != 1:
                raise SourceDataError("PowerPoint conversion did not produce exactly one PDF")
            prefix = png_dir / "slide"
            subprocess.run(
                [pdftoppm, "-png", "-r", "110", str(pdf_files[0]), str(prefix)],
                check=True,
                capture_output=True,
                text=True,
                timeout=RENDER_TIMEOUT_SECONDS,
            )
        except subprocess.TimeoutExpired as exc:
            raise SourceCapabilityError("PPTX rendering exceeded the verifier timeout") from exc
        except subprocess.CalledProcessError as exc:
            detail = (exc.stderr or exc.stdout or "conversion failed").strip()
            raise SourceDataError(f"PowerPoint rendering failed: {detail[:300]}") from exc

        images = sorted(png_dir.glob("slide-*.png"), key=slide_sort_key)
        if not images:
            raise SourceDataError("PowerPoint rendering produced no slide images")
        slide_count = len(images)
        selected: list[RenderedSlide] = []
        total_bytes = 0
        truncated = slide_count > MAX_RENDERED_SLIDES
        for image in images[:MAX_RENDERED_SLIDES]:
            raw = image.read_bytes()
            if total_bytes + len(raw) > MAX_RENDERED_BYTES:
                truncated = True
                break
            width, height = png_dimensions(raw)
            selected.append(RenderedSlide(
                number=slide_number_from_path(image),
                width=width,
                height=height,
                image_base64=base64.b64encode(raw).decode("ascii"),
            ))
            total_bytes += len(raw)
        if not selected:
            raise SourceDataError("PowerPoint rendering exceeded the configured image evidence limit")
        return PptxRenderSlidesOutput(
            text=text,
            slide_count=slide_count,
            slides=selected,
            truncated=truncated,
            renderer="libreoffice-impress+pdftoppm",
        )


def slide_is_hidden(slide: Any) -> bool:
    """Whether PowerPoint would skip this slide when presenting.

    A hidden slide is stored as ``<p:sld show="0">``. python-pptx exposes no
    accessor for it, so this reads the underlying element defensively: a future
    python-pptx that reshapes the element must degrade to "visible" rather than
    raise, since ``AttributeError`` is absent from ``agent_failure_types`` and
    would fail the whole run instead of one check.

    Args:
        slide: python-pptx slide object.

    Returns:
        True when the slide carries an explicit ``show="0"``.
    """
    element = getattr(slide, "_element", None)
    getter = getattr(element, "get", None)
    if not callable(getter):
        return False
    return str(getter("show") or "").strip() in {"0", "false"}


def slide_sort_key(path: Path) -> tuple[int, str]:
    """Sort Poppler's ``slide-12.png`` names numerically."""
    return slide_number_from_path(path), path.name


def slide_number_from_path(path: Path) -> int:
    """Read Poppler's trailing page number, with a stable fallback."""
    stem = path.stem.rsplit("-", 1)[-1]
    try:
        return int(stem)
    except ValueError:
        return 0


def png_dimensions(raw: bytes) -> tuple[int, int]:
    """Read PNG IHDR dimensions without adding an image-processing dependency."""
    if len(raw) < 24 or raw[:8] != b"\x89PNG\r\n\x1a\n" or raw[12:16] != b"IHDR":
        raise SourceDataError("PowerPoint renderer returned an invalid PNG")
    width = int.from_bytes(raw[16:20], "big")
    height = int.from_bytes(raw[20:24], "big")
    if width <= 0 or height <= 0:
        raise SourceDataError("PowerPoint renderer returned a PNG with invalid dimensions")
    return width, height


def iter_slide_text(slide: Any) -> list[str]:
    """Collects text from a slide's text frames and tables.

    Args:
        slide: python-pptx slide object.

    Returns:
        Non-empty text fragments from the slide.
    """
    fragments: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                fragments.append(text)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    fragments.append(" | ".join(cells))
    return fragments


def append_limited(parts: list[str], text: str, total: int, limit: int) -> int:
    """Appends text and returns the updated character count.

    Args:
        parts: Text fragments collected so far.
        text: Text fragment to append.
        total: Current approximate character count.
        limit: Maximum desired character count.

    Returns:
        Updated approximate character count.
    """
    parts.append(text)
    return total + len(text) + 1


COMMANDS = (ExtractText(), RenderSlides(), InspectDeck())
